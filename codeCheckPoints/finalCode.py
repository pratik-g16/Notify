from flask import Flask, flash, request, redirect, url_for, render_template
import urllib.request
import os
from werkzeug.utils import secure_filename
import numpy as np
import cv2
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import moviepy.editor as mp
import speech_recognition as sr 
import math 
from fpdf import FPDF
from gtts import gTTS
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from spellchecker import SpellChecker 
import pytesseract
pytesseract.pytesseract.tesseract_cmd = 'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'
import pptx
from pptx import Presentation 
from pptx.util import Inches   
from pydub import AudioSegment
file_name = 'notes'

curr, prev = 0, 0
a = []
num = 0


def convertToGray(frame, no):
    image1 = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
    mask = np.full(image1.shape, 255)
    mod = image1-mask
    mod = mod*-1
    mod = mod.astype(np.uint8)
    lower = np.array([0])
    upper = np.array([50])
    mask1 = cv2.inRange(mod, lower, upper)
    final = mask1
    # counting the number of pixels
    white = np.sum(final == 255)
    black = np.sum(final == 0)
    curr = (black/(white+black))
    curr = round(curr*100, 3)
    a.append(curr)
    return final, curr


def keyFrame(video):
    keyframes = []
    videoFile = video
    cap = cv2.VideoCapture(videoFile)
    frameRate = cap.get(5)  # frame rate
    page = 0
    key = 0
    aud = []
    while(cap.isOpened()):
        frameId = cap.get(1)  # current frame number
        ret, frame = cap.read()
        if (ret != True):
            break
        if (frameId % math.floor(frameRate) == 0):
            page += 1
            frame, pct = convertToGray(frame, int(page))
            if(page == 1):
                pprev = [pct, frame]
                prev = [pct, frame]
                curr = [pct, frame]
            if(page == 2):
                pprev = [prev[0], prev[1]]
                prev = [pct, frame]
                curr = [pct, frame]
            else:
                pprev = [prev[0], prev[1]]
                prev = [curr[0], curr[1]]
                curr = [pct, frame]
                # if(prev[0] > curr[0]*1.1 and pprev[0] <= prev[0]):
                if(prev[0] > 1.1* curr[0] and pprev[0] < 1.1* prev[0]):    
                    aud.append(page-1)
                    key += 1
                    keyframes.append(prev[1])    
    if(round(prev[0],1)<=round(curr[0],1)*1.1):
    # if(prev[0] < curr[0] and pprev[0] <= prev[0]):
        keyframes.append(curr[1])
        aud.append(page-1)
    cap.release()
    return aud, keyframes


 
# def keyFrame(video):
#     keyframes = []
#     videoFile = video
#     cap = cv2.VideoCapture(videoFile)
#     frameRate = cap.get(5)  # frame rate
#     page = 0
#     key = 0
#     aud = []
#     while(cap.isOpened()):
#         frameId = cap.get(1)  # current frame number
#         ret, frame = cap.read()
#         if (ret != True):
#             # cv2.imwrite(filename, frame)
#             break
#         if (frameId % math.floor(frameRate) == 0):
#             page += 1 
#             frame, pct = convertToGray(frame, int(page))
#             if(page == 1):
#                 prev = [pct, frame]
#                 curr = [pct, frame]
#                 # print("page==1", curr, prev)
#             else:
#                 prev = [curr[0], curr[1]]
#                 curr = [pct, frame]
#                 # print('[[[', curr, prev)
#                 if(prev[0] > curr[0]*1.1):
#                     # print("----", page-1)
#                     key += 1 
#                     keyframes.append(prev[1])

#     if(prev[0] <= curr[0]): 
#         keyframes.append(curr[1])
#     cap.release()
#     print("Done!")
#     # print(a)
#     return aud, keyframes

def chunkAudio(a):
    aud=[]
    k=0 
    text=[]
    


    aud.append([0,a[0]])
    song = AudioSegment.from_mp3('static\\output\\output.wav')
    extract = song[0:a[0]*1000]
    k+=1
    # # Saving
    extract.export("static\\output\\notes-extract.wav", format="wav")
    r = sr.Recognizer()
    with sr.AudioFile("static\\output\\notes-extract.wav") as source:
            audio = r.record(source) 
            s = r.recognize_google(audio)
            # text.append(s)
            s = s.split("\n")
            p = []
            j = 0
            for i in s:
                i = i.strip()
                if len(i) >= 70:
                    for k in range(len(i)//70):
                        p.append(i[j:j+70])
                        j += 70
                    p.append(i[j:])
                else:
                    p.append(i)
            # print(p)
            s="\n".join(p)
            text.append(s)



    for i in range(1,len(a)):
        aud.append([a[i-1],a[i]])
        # song = AudioSegment.from_mp3('C:\\Users\\HP\\Desktop\\project\\static\\output\\output.wav')
        song = AudioSegment.from_mp3('static\\output\\output.wav')
        extract = song[a[i-1]*1000:a[i]*1000]
        k+=1
        # # Saving
        # extract.export(file_name+'-extract.wav', format="wav")
        extract.export("static\\output\\notes-extract.wav", format="wav")
        r = sr.Recognizer()
        # with sr.AudioFile("C:\\Users\\HP\\Desktop\\project\\notes-extract.wav") as source:
        with sr.AudioFile("static\\output\\notes-extract.wav") as source:
            audio = r.record(source) 
            s = r.recognize_google(audio)
            # text.append(s) 
            s = s.split("\n")
            p = []
            j = 0
            for i in s:
                i = i.strip()
                if len(i) >= 70:
                    for k in range(len(i)//70):
                        p.append(i[j:j+70])
                        j += 70
                    p.append(i[j:])
                else:
                    p.append(i)
            # print(p)
            s="\n".join(p)
            text.append(s)
    return text



def segment(image):
    # grayscale
    lines = []
    # gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    gray = image
    # binary
    ret, thresh = cv2.threshold(gray, 127, 255, cv2.THRESH_BINARY_INV)

    # dilation
    kernel = np.ones((5, 100), np.uint8)
    img_dilation = cv2.dilate(thresh, kernel, iterations=1)
    # find contours
    ctrs, hier = cv2.findContours(
        img_dilation.copy(), cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    # sort contours
    sorted_ctrs = sorted(ctrs, key=lambda ctr: cv2.boundingRect(ctr)[1])

    for i, ctr in enumerate(sorted_ctrs):
        # Get bounding box
        x, y, w, h = cv2.boundingRect(ctr)
        # Getting ROI
        roi = image[y:y+h, x:x+w]
        lines.append(roi)
    return lines




def getTranscript(video):
   video = mp.VideoFileClip(video)
   video.audio.write_audiofile(r"static\\output\\output.wav") 
   r = sr.Recognizer()
   with sr.AudioFile("static\\output\\output.wav") as source:
      audio = r.record(source) 
   s = r.recognize_google(audio)
   return s


# def getTranscript(video):
#     text=''
#     video = mp.VideoFileClip(video)
#     video.audio.write_audiofile(r"static\\output\\output.wav") 
#     # r = sr.Recognizer()
#     # with sr.AudioFile("static\\output\\output.wav") as source:
#     #     audio = r.record(source) 
#     # s = r.recognize_google(audio)
#     song = AudioSegment.from_mp3('static\\output\\output.wav')
#     extract1 = song[:300000]
#     # # Saving
#     extract1.export("static\\output\\notes-extract.wav", format="wav")
#     r = sr.Recognizer()
#     with sr.AudioFile("static\\output\\notes-extract.wav") as source:
#             audio = r.record(source) 
#             text = r.recognize_google(audio)
#     extract2 = song[300001:]
#     # # Saving
#     extract2.export("static\\output\\notes-extract.wav", format="wav")
#     r = sr.Recognizer()
#     with sr.AudioFile("static\\output\\notes-extract.wav") as source:
#             audio = r.record(source) 
#             text +=' '+ r.recognize_google(audio)
#     return text

def getTranscriptForPPT(video):
    text=''
    video = mp.VideoFileClip(video)
    video.audio.write_audiofile(r"static\\output\\output.wav") 
    # r = sr.Recognizer()
    # with sr.AudioFile("static\\output\\output.wav") as source:
    #     audio = r.record(source) 
    # s = r.recognize_google(audio)
    song = AudioSegment.from_mp3('static\\output\\output.wav')
    extract1 = song[:200000]
    # # Saving
    extract1.export("static\\output\\notes-extract.wav", format="wav")
    r = sr.Recognizer()
    with sr.AudioFile("static\\output\\notes-extract.wav") as source:
            audio = r.record(source) 
            text += r.recognize_google(audio) 
    extract2 = song[200001:]
    # # Saving
    extract2.export("static\\output\\notes-extract.wav", format="wav")
    r = sr.Recognizer()
    with sr.AudioFile("static\\output\\notes-extract.wav") as source:
            audio = r.record(source) 
            text +=' '+ r.recognize_google(audio) 
    return text

def summarize(text):
    stopWords = set(stopwords.words("english"))
    words = word_tokenize(text)
    # print(words)
    # Creating a frequency table to keep the
    # score of each word

    freqTable = dict()
    for word in words:
        word = word.lower()
        if word in stopWords:
            continue
        if word in freqTable:
            freqTable[word] += 1
        else:
            freqTable[word] = 1
    # print(freqTable)
    sentences = sent_tokenize(text)
    # print(sentences)
    sentenceValue = dict()
    sumValues = 0
    for sentence in sentences:
        for word, freq in freqTable.items():
            if word in sentence.lower():
                if sentence in sentenceValue:
                    sentenceValue[sentence] += freq
                else:
                    sentenceValue[sentence] = freq
    for sentence in sentenceValue:
        sumValues += sentenceValue[sentence]
    # print(sentenceValue)
    # print(sumValues)

    average = int(sumValues / len(sentenceValue))
    # print(average)
    # Storing sentences into our summary.
    summary = ''
    for sentence in sentences:
        if (sentence in sentenceValue) and (sentenceValue[sentence] > (1.2 * average)):
            summary += " " + sentence
    return summary


def convertToPdf(text, name):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Times", size=20)
    a = text
    s = a.split("\n")
    p = []
    j = 0
    for i in s:
        i = i.strip()
        if len(i) >= 60:
            for k in range(len(i)//60):
                p.append(i[j:j+60])
                j += 60
            p.append(i[j:])
        else:
            p.append(i)
    # p.remove('')

    for i in p:
        pdf.cell(200, 10, txt=i, ln=1, align='L')
    pdf.output(name+'.pdf')


def getAudio(text):
    mytext = ' '.join(text.split('\n'))
    language = 'en'
    myobj = gTTS(text=mytext, lang=language, slow=False)
    myobj.save("static\\output\\audioNotes.mp3")

def getText(img): 
    # img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    # cv2.imshow('img', img)
    # cv2.waitKey()
    text = pytesseract.image_to_string(img)
    # print(type(text))
    text = text[:len(text)-1]
    # print(text)
    return text

def sendMail():
    fromaddr = "notes.notifier@gmail.com"
    toaddr = "guptapn@rknec.edu,dholwanimona@gmail.com,varmans_1@rknec.edu,guptava@rknec.edu,deshmukhsp_1@rknec.edu"
    msg = MIMEMultipart() 
    msg['From'] = fromaddr 
    msg['To'] = toaddr 
    msg['Subject'] = "Notes for today's class" 
    body = "PFA notes for today's lecture." 
    msg.attach(MIMEText(body, 'plain')) 
    filename = "Notes.zip"
    # os.chdir("static\\output")
    attachment = open("Notes.zip", "rb") 
    p = MIMEBase('application', 'octet-stream') 
    p.set_payload((attachment).read()) 
    encoders.encode_base64(p) 
    p.add_header('Content-Disposition', "attachment; filename= %s" % filename) 
    msg.attach(p) 
    s = smtplib.SMTP('smtp.gmail.com', 587) 
    s.starttls() 
    s.login(fromaddr, "MakeItHappen") 
    text = msg.as_string() 
    s.sendmail(fromaddr, toaddr.split(","), text) 
    s.quit()
    # print("done1")


def convertToPpt(videoFile):  
    cap = cv2.VideoCapture(videoFile)
    frameRate = cap.get(5)  # frame rate
    page = 0
    key = 0 
    keyframes=[]
    while(cap.isOpened()):
        frameId = cap.get(1)  # current frame number
        ret, frame = cap.read()
        if (ret != True):
            # cv2.imwrite(filename, frame)
            break
        if (frameId % math.floor(frameRate) == 0):
            page += 1 
            _,pct= convertToGray(frame, int(page))
            if(page == 1):
                prev = [pct, frame]
                curr = [pct, frame]
                # print("page==1", curr, prev)
            else: 
                prev = [curr[0], curr[1]]
                curr = [pct, frame]
                # print('[[[', curr, prev)
                # if(prev[0] > curr[0]*1.1):
                #     print("----", page-1)
                #     key += 1
                #     # cv2.imwrite(filename, prev[1])
                #     keyframes.append(prev[1])
                # print(abs(prev[0]-curr[0]))
                if(abs(prev[0]-curr[0])>=0.1): 
                    keyframes.append(prev[1])
                    

    # if(abs(prev[0]-curr[0])>=0.1):
    # if(prev[0] != curr[0]):
        # print("----", page-1)
        # cv2.imwrite(filename, curr[1])
    keyframes.append(curr[1])
    cap.release()
    print("Done!")
    # print(a)

    # x = np.arange(1, len(a)+1, 1)
    # plt.scatter(x, np.array(a), color='b', s=100)
    # plt.ylim(0, 10)
    # plt.xticks(x)
    # plt.xlabel('Pages')
    # plt.ylabel('% of text')
    # plt.show()


    # for i in keyframes:
    #     cv2.imshow('a',i)
    #     cv2.waitKey()

    # t=''

    # for i in keyframes:
        # img=i
        # # img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        # cv2.imshow('img', img)
        # cv2.waitKey()
    prs = Presentation()  
    for i in keyframes:
        cv2.imwrite('static\\output\\pptframe.jpg',i)
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = Inches(1)
        # pic = slide.shapes.add_picture(img_path, left, top) 
        pic = slide.shapes.add_picture('static\\output\\pptframe.jpg', pptx.util.Inches(0), pptx.util.Inches(0.7),
                                        width=pptx.util.Inches(10), height=pptx.util.Inches(6))
    
    prs.save('static\\output\\notes.pptx')





def getOutputForHandwritten(video):
    # segmentedLines = []

    audioKeyFrames, keyframes = keyFrame(video)
    # print(audioKeyFrames)
    # print(audioKeyFrames)
    
    # for i in keyframes:
    #     cv2.imshow('frame', i)
    #     cv2.waitKey()
    #     frameLines = segment(i)
    #     segmentedLines += frameLines

    # cv2.imwrite('l.jpg', segmentedLines[1])
    # for i in segmentedLines:
    #     cv2.imshow('l', i)
    #     cv2.waitKey()
        # handwriting recognition

    # lines = segment(image)
    # for i in lines:
    #     cv2.imshow('a', i)
    #     cv2.waitKey()

    # Input text - to summarize
    text = ''
    for i in keyframes:
        text+=getText (i) 
    # print(text)
    summary = summarize(text)

    convertToPdf(text, 'static\\output\\notes')
    convertToPdf(summary, 'static\\output\\summary')
    getAudio(text)
    text=getTranscript(video) 
    convertToPdf(text,'static\\output\\transcripts')
    # print(audioKeyFrames)
    t=chunkAudio(audioKeyFrames)
    prs = Presentation()  
    for i in range(len(keyframes)):
        cv2.imwrite('static\\output\\pptframe.jpg',keyframes[i])
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = Inches(1)
        # pic = slide.shapes.add_picture(img_path, left, top) 
        pic = slide.shapes.add_picture('static\\output\\pptframe.jpg', pptx.util.Inches(0), pptx.util.Inches(0.7),
                                        width=pptx.util.Inches(10), height=pptx.util.Inches(6))

        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1) 
        txBox = slide.shapes.add_textbox(left, top,
                                                width, height)
        tf = txBox.text_frame
        tf.text = t[i]
        p = tf.add_paragraph()
    
    prs.save('static\\output\\summary.pptx')
    
    
    
    



 






def getOutputForPpt(video):
    convertToPpt(video) 
    text=getTranscriptForPPT(video) 
    # print(text)
    convertToPdf(text,'static\\output\\transcripts')
    # if(os.path.exists(os.path.join('static/output/', 'pptframe.jpg'))):
    #     os.remove(os.path.join('static/output/', 'pptframe.jpg'))




from zipfile import ZipFile
def makeZip(filenames):
    os.chdir("static\\output")
    zipObj = ZipFile('Notes.zip', 'w')
    # print(filenames)
    for i in filenames:
        zipObj.write(i)
    zipObj.close()


 





















import flask

app = Flask(__name__)

UPLOAD_FOLDER = 'static/uploads/'

app.secret_key = "secret key"
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
# app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

ALLOWED_EXTENSIONS = set(['mp4'])


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route('/')
def home():
    if(os.path.exists(os.path.join(app.config['UPLOAD_FOLDER'], 'video.mp4'))):
        os.remove(os.path.join(app.config['UPLOAD_FOLDER'], 'video.mp4'))
    
    if(os.path.exists(os.path.join('static/output/', 'notes.pdf'))):
        os.remove(os.path.join('static/output/', 'notes.pdf'))

    if(os.path.exists(os.path.join('static/output/', 'summary.pdf'))):
        os.remove(os.path.join('static/output/', 'summary.pdf'))

    if(os.path.exists(os.path.join('static/output/', 'transcripts.pdf'))):
        os.remove(os.path.join('static/output/', 'transcripts.pdf'))

    if(os.path.exists(os.path.join('static/output/', 'output.wav'))):
        os.remove(os.path.join('static/output/', 'output.wav'))

    if(os.path.exists(os.path.join('static/output/', 'audioNotes.mp3'))):
        os.remove(os.path.join('static/output/', 'audioNotes.mp3'))

    if(os.path.exists(os.path.join('static/output/', 'notes-extract.wav'))):
        os.remove(os.path.join('static/output/', 'notes-extract.wav'))

    if(os.path.exists(os.path.join('static/output/', 'summary.pptx'))):
        os.remove(os.path.join('static/output/', 'summary.pptx'))

    if(os.path.exists(os.path.join('static/output/', 'notes.pptx'))):
        os.remove(os.path.join('static/output/', 'notes.pptx'))

    if(os.path.exists(os.path.join('static/output/', 'pptframe.jpg'))):
        os.remove(os.path.join('static/output/', 'pptframe.jpg'))

    if(os.path.exists(os.path.join('static/output/', 'Notes.zip'))):
        os.remove(os.path.join('static/output/', 'Notes.zip'))

    return render_template('index.html')



@app.route('/uploader', methods=['POST'])
def upload_image2(): 
    if 'file' not in request.files: 
        return redirect(request.url) 
    file = request.files['file']
    if file.filename == '': 
        return redirect(request.url)
    if file and allowed_file(file.filename): 
        filename = secure_filename("video.mp4") 
        # print(filename)
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename)) 
        option = request.form['options'] 
        # print(option) 
        if(option=='hw'):
            getOutputForHandwritten('static\\uploads\\video.mp4')
            makeZip(["notes.pdf","audioNotes.mp3","transcripts.pdf","summary.pptx"])
            # print("done1")
            # sendMail()
        else:
            getOutputForPpt('static\\uploads\\video.mp4')
            # makeZip(["summary.pptx"])  
            makeZip(["notes.pptx","transcripts.pdf"])  
            # print("done1")
            # sendMail()
        # return "done"
        # sendMail()
        return flask.send_file("static\\output\\Notes.zip", as_attachment=True)
        return  flask.redirect("/")
        # return "done"
    else: 
        return redirect(request.url)


		
if __name__ == '__main__':
   app.run(debug = True)






