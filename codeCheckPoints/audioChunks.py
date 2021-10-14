# import os
# from zipfile import ZipFile

# os.chdir('static\\output')
# zipObj = ZipFile('output.zip', 'w')
# zipObj.write("notes.pptx")
# zipObj.close()


# from pydub import AudioSegment
# from pydub.utils import make_chunks

# myaudio = AudioSegment.from_file("C:\\Users\\HP\\Desktop\\project\\static\\output\\output.wav" , "wav") 
# chunk_length_ms = 1000 # pydub calculates in millisec
# chunks = make_chunks(myaudio, chunk_length_ms) #Make chunks of one sec

# #Export all of the individual chunks as wav files

# for i, chunk in enumerate(chunks):
#     chunk_name = "chunk{0}.wav".format(i)
#     print ("exporting", chunk_name)
#     chunk.export(chunk_name, format="wav")

# import speech_recognition as sr
# from pydub import AudioSegment
# file_name = 'notes'
# # k=0
# # song = AudioSegment.from_mp3('output.wav')
# # extract = song[1:50000]
# # k+=1
# #     # # Saving
# # extract.export(file_name+str(k)+'-extract.wav', format="wav")
    



# a=[5,10,14,25,30]
# aud=[]
# k=0 

# for i in range(1,len(a)):
#     aud.append([a[i-1],a[i]])
#     song = AudioSegment.from_mp3('output.wav')
#     extract = song[a[i-1]*1000:a[i]*1000]
#     k+=1
#     # # Saving
#     extract.export(file_name+str(k)+'-extract.wav', format="wav")
    

# extract = song[a[i]*1000:]
# k+=1
#     # # Saving
# extract.export(file_name+str(k)+'-extract.wav', format="wav")
# def getTranscript(): 
#    r = sr.Recognizer()
#    with sr.AudioFile("C:\\Users\\HP\\Desktop\\project\\notes5-extract.wav") as source:
#       audio = r.record(source) 
#    s = r.recognize_google(audio)
#    return s

# print(getTranscript())






 # import required things
from pptx import Presentation 
from pptx.util import Inches, Pt
import speech_recognition as sr
from pydub import AudioSegment
file_name = 'notes'
a=[78,139,273]



aud=[]
k=0 
text=[]
ppt = Presentation() 


aud.append([0,a[0]])
song = AudioSegment.from_mp3('C:\\Users\\HP\\Desktop\\project\\static\\output\\output.wav')
extract = song[0:a[0]*1000]
k+=1
# # Saving
extract.export(file_name+'-extract.wav', format="wav")
r = sr.Recognizer()
with sr.AudioFile("C:\\Users\\HP\\Desktop\\project\\notes-extract.wav") as source:
        audio = r.record(source) 
        s = r.recognize_google(audio)
        # text.append(s)
        text=''
        s = s.split("\n")
        p = []
        j = 0
        for i in s:
            i = i.strip()
            if len(i) >= 70:
                for k in range(len(i)//70):
                    p.append(i[j:j+70])
                    j += 70
                p.append(i[j:])
            else:
                p.append(i)
        # print(p)
        s="\n".join(p)
        print(s)
        blank_slide_layout = ppt.slide_layouts[6] 
        slide = ppt.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1) 
        txBox = slide.shapes.add_textbox(left, top,
                                            width, height)
        tf = txBox.text_frame
        tf.text = s
        p = tf.add_paragraph() 



for i in range(1,len(a)):
    aud.append([a[i-1],a[i]])
    song = AudioSegment.from_mp3('C:\\Users\\HP\\Desktop\\project\\static\\output\\output.wav')
    extract = song[a[i-1]*1000:a[i]*1000]
    k+=1
    # # Saving
    extract.export(file_name+'-extract.wav', format="wav")
    r = sr.Recognizer()
    with sr.AudioFile("C:\\Users\\HP\\Desktop\\project\\notes-extract.wav") as source:
        audio = r.record(source) 
        s = r.recognize_google(audio)
        # text.append(s)
        text=''
        s = s.split("\n")
        p = []
        j = 0
        for i in s:
            i = i.strip()
            if len(i) >= 70:
                for k in range(len(i)//70):
                    p.append(i[j:j+70])
                    j += 70
                p.append(i[j:])
            else:
                p.append(i)
        # print(p)
        s="\n".join(p) 
        print(s)
        blank_slide_layout = ppt.slide_layouts[6] 
        slide = ppt.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1) 
        txBox = slide.shapes.add_textbox(left, top,
                                            width, height)
        tf = txBox.text_frame
        tf.text = s
        p = tf.add_paragraph() 

ppt.save('test_2.pptx')

print("done")


# for i in text:
#     print(i)
#     print("-------") 
#     # Creating Object
    

# # import required things
# from pptx import Presentation 
# from pptx.util import Inches, Pt
  
# ppt = Presentation() 
# blank_slide_layout = ppt.slide_layouts[6] 
# slide = ppt.slides.add_slide(blank_slide_layout)
# left = top = width = height = Inches(1) 
# txBox = slide.shapes.add_textbox(left, top,
#                                     width, height)
# tf = txBox.text_frame
# tf.text = text[i]
# p = tf.add_paragraph() 

# ppt.save('test_2.pptx')

# print("done")